\
from __future__ import annotations

import io
from typing import Tuple, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.dml import MSO_FILL
from pptx.util import Pt


def _set_slide_background(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _is_title(shape) -> bool:
    if not shape.has_text_frame:
        return False
    if not getattr(shape, "is_placeholder", False):
        # heuristics: large text box at top might still be title; we keep it simple
        return False
    phf = shape.placeholder_format
    return phf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)


def apply_style(pptx_bytes: bytes, style: str) -> bytes:
    """
    style:
      - "dark_neon": dark background, orange titles, white body, yellow emphasis
      - "clean_minimal": white background, dark titles, gray body, subtle accent
    """
    prs = Presentation(io.BytesIO(pptx_bytes))

    if style == "dark_neon":
        bg = RGBColor(10, 10, 10)
        title_color = RGBColor(255, 140, 0)   # orange
        body_color = RGBColor(245, 245, 245)  # near white
    elif style == "clean_minimal":
        bg = RGBColor(255, 255, 255)
        title_color = RGBColor(20, 20, 20)
        body_color = RGBColor(60, 60, 60)
    else:
        raise ValueError(f"Unknown style: {style}")

    for slide in prs.slides:
        _set_slide_background(slide, bg)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            is_title = _is_title(shape)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font
                    font.name = "Calibri"  # safe default; you can replace with Noto Sans TC if installed
                    if is_title:
                        font.bold = True
                        font.size = font.size or Pt(36)
                        font.color.rgb = title_color
                    else:
                        font.size = font.size or Pt(20)
                        font.color.rgb = body_color

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
